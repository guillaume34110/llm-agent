"""Office tools: Excel, Word, PowerPoint, advanced PDF — niveau grand-maître bureautique."""
from __future__ import annotations
import json
from pathlib import Path
from typing import Any

from monkey.tools.files import _resolve


def _ok(m: str) -> str: return f"OK: {m}"
def _err(m: str) -> str: return f"ERREUR: {m}"


# ─── Excel (.xlsx) ───────────────────────────────────────────────────────────

def xlsx_create(path: str, sheets: dict | list | None = None) -> str:
    """Create an .xlsx file. sheets: {sheet_name: [[row1], [row2]...]} or list of rows for default sheet."""
    try:
        import openpyxl
    except ImportError:
        return _err("openpyxl manquant")
    wb = openpyxl.Workbook()
    if isinstance(sheets, list):
        sheets = {"Sheet1": sheets}
    if sheets:
        wb.remove(wb.active)
        for name, rows in sheets.items():
            ws = wb.create_sheet(name)
            for row in rows or []:
                ws.append(row)
    p = _resolve(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    wb.save(p)
    return _ok(f"xlsx → {p} ({len(sheets or {})} feuille(s))")


def xlsx_read(path: str, sheet: str = "", max_rows: int = 1000) -> str:
    """Read xlsx → JSON {sheet_name: [[cells]]}. If sheet given, only that one."""
    try:
        import openpyxl
    except ImportError:
        return _err("openpyxl manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    wb = openpyxl.load_workbook(p, data_only=True, read_only=True)
    out: dict[str, list] = {}
    sheets = [sheet] if sheet else wb.sheetnames
    for s in sheets:
        if s not in wb.sheetnames:
            continue
        ws = wb[s]
        rows = []
        for i, r in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows:
                break
            rows.append([c for c in r])
        out[s] = rows
    return json.dumps(out, ensure_ascii=False, default=str, indent=2)[:12000]


def xlsx_write_cells(path: str, sheet: str, cells: dict) -> str:
    """Update specific cells. cells: {"A1": "value", "B2": 42, "C3": "=SUM(A1:A10)"}."""
    try:
        import openpyxl
    except ImportError:
        return _err("openpyxl manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    wb = openpyxl.load_workbook(p)
    if sheet not in wb.sheetnames:
        wb.create_sheet(sheet)
    ws = wb[sheet]
    for ref, val in (cells or {}).items():
        ws[ref] = val
    wb.save(p)
    return _ok(f"xlsx écrit {len(cells or {})} cellule(s) dans {sheet}")


def xlsx_append_rows(path: str, sheet: str, rows: list) -> str:
    """Append rows to a sheet (creates file/sheet if needed)."""
    try:
        import openpyxl
    except ImportError:
        return _err("openpyxl manquant")
    p = _resolve(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    if p.exists():
        wb = openpyxl.load_workbook(p)
    else:
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
    if sheet not in wb.sheetnames:
        wb.create_sheet(sheet)
    ws = wb[sheet]
    for r in rows or []:
        ws.append(r)
    wb.save(p)
    return _ok(f"xlsx +{len(rows or [])} lignes dans {sheet}")


def xlsx_to_csv(path: str, output_path: str, sheet: str = "") -> str:
    """Export a sheet to CSV."""
    try:
        import openpyxl, csv
    except ImportError:
        return _err("openpyxl manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    wb = openpyxl.load_workbook(p, data_only=True, read_only=True)
    s = sheet or wb.sheetnames[0]
    if s not in wb.sheetnames:
        return _err(f"feuille inconnue: {s}")
    ws = wb[s]
    out = _resolve(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        for r in ws.iter_rows(values_only=True):
            w.writerow(r)
    return _ok(f"CSV → {out}")


# ─── Word (.docx) ────────────────────────────────────────────────────────────

def docx_create(path: str, title: str = "", paragraphs: list | None = None,
                content: str = "") -> str:
    """Create a .docx. Use paragraphs list or markdown-ish `content` (# h1 / ## h2 / blank line = new para)."""
    try:
        from docx import Document
    except ImportError:
        return _err("python-docx manquant")
    doc = Document()
    if title:
        doc.add_heading(title, level=0)
    if paragraphs:
        for para in paragraphs:
            if isinstance(para, dict):
                style = para.get("style", "")
                text = para.get("text", "")
                if style.startswith("h") and style[1:].isdigit():
                    doc.add_heading(text, level=int(style[1:]))
                elif style == "list":
                    doc.add_paragraph(text, style="List Bullet")
                else:
                    doc.add_paragraph(text)
            else:
                doc.add_paragraph(str(para))
    elif content:
        for line in content.split("\n"):
            s = line.rstrip()
            if not s:
                continue
            if s.startswith("# "):
                doc.add_heading(s[2:], level=1)
            elif s.startswith("## "):
                doc.add_heading(s[3:], level=2)
            elif s.startswith("### "):
                doc.add_heading(s[4:], level=3)
            elif s.startswith(("- ", "* ")):
                doc.add_paragraph(s[2:], style="List Bullet")
            else:
                doc.add_paragraph(s)
    p = _resolve(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    doc.save(p)
    return _ok(f"docx → {p}")


def docx_read(path: str, max_chars: int = 12000) -> str:
    """Extract text from a .docx file."""
    try:
        from docx import Document
    except ImportError:
        return _err("python-docx manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    doc = Document(p)
    parts = [para.text for para in doc.paragraphs if para.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)[:max_chars]


def docx_replace(path: str, replacements: dict) -> str:
    """Replace placeholders in a .docx (keeps formatting). replacements: {"{{name}}": "Jean"}."""
    try:
        from docx import Document
    except ImportError:
        return _err("python-docx manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    doc = Document(p)
    n = 0
    for para in doc.paragraphs:
        for run in para.runs:
            for old, new in (replacements or {}).items():
                if old in run.text:
                    run.text = run.text.replace(old, str(new))
                    n += 1
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for old, new in (replacements or {}).items():
                    if old in cell.text:
                        cell.text = cell.text.replace(old, str(new))
                        n += 1
    doc.save(p)
    return _ok(f"docx {n} remplacement(s) dans {p}")


# ─── PowerPoint (.pptx) ──────────────────────────────────────────────────────

def pptx_create(path: str, slides: list) -> str:
    """Create a .pptx. slides: [{"title": "...", "content": "..." or ["bullet1", "bullet2"]}]."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return _err("python-pptx manquant")
    prs = Presentation()
    for sl in slides or []:
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = sl.get("title", "")
        body = slide.placeholders[1].text_frame
        content = sl.get("content", "")
        if isinstance(content, list):
            body.text = content[0] if content else ""
            for line in content[1:]:
                p = body.add_paragraph()
                p.text = line
        else:
            body.text = content
    p = _resolve(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(p)
    return _ok(f"pptx → {p} ({len(slides or [])} slides)")


def pptx_read(path: str) -> str:
    """Extract text from a .pptx (titles + content per slide)."""
    try:
        from pptx import Presentation
    except ImportError:
        return _err("python-pptx manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    prs = Presentation(p)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        parts.append(txt)
    return "\n".join(parts)[:12000]


# ─── PDF avancé ──────────────────────────────────────────────────────────────

_PDF_PAGE_CHAR_BUDGET = 14000


def pdf_extract_text(path: str, pages: str = "") -> str:
    """Extract text from a PDF. pages: '1-3,5' or empty for all. OCR fallback on scanned pages."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return _err("pypdf manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    try:
        reader = PdfReader(str(p))
    except Exception as e:
        return _err(f"PDF illisible ({type(e).__name__}: {e})")
    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")
        except Exception:
            return _err("PDF chiffré — fournir mot de passe via pdf_decrypt")
    total = len(reader.pages)
    if total == 0:
        return _err("PDF vide (0 pages)")
    targets = list(_parse_pages(pages, total)) if pages else list(range(total))
    parts: list[str] = [f"PDF: {p.name} ({total} pages, extract {len(targets)})"]
    ocr_pages: list[int] = []
    used_chars = 0
    truncated_at: int | None = None
    for i in targets:
        if not (0 <= i < total):
            continue
        try:
            text = reader.pages[i].extract_text() or ""
        except Exception as ex:
            text = f"[extract error: {ex}]"
        text = text.strip()
        if len(text) < 20:
            ocr = _ocr_pdf_page(p, i)
            if ocr:
                text = ocr
                ocr_pages.append(i + 1)
            elif not text:
                text = "[no extractable text — likely scanned/image page; OCR unavailable]"
        block = f"--- Page {i+1}/{total} ---\n{text}"
        if used_chars + len(block) > _PDF_PAGE_CHAR_BUDGET:
            remaining = [t + 1 for t in targets[targets.index(i):]]
            truncated_at = i + 1
            parts.append(
                f"[TRUNCATED at page {truncated_at} — budget {_PDF_PAGE_CHAR_BUDGET} chars reached. "
                f"Remaining pages {remaining[0]}-{remaining[-1]}. "
                f"Call pdf_extract_text(path, pages='{remaining[0]}-{remaining[-1]}') to continue.]"
            )
            break
        parts.append(block)
        used_chars += len(block)
    if ocr_pages:
        parts.insert(1, f"[OCR fallback applied on pages: {','.join(str(x) for x in ocr_pages)}]")
    return "\n".join(parts)


def _ocr_pdf_page(pdf_path: Path, page_idx: int) -> str:
    """Rasterize one PDF page and OCR it. Returns empty string if OCR unavailable."""
    try:
        from pdf2image import convert_from_path
        import pytesseract
        import subprocess
        if subprocess.run(["which", "tesseract"], capture_output=True).returncode != 0:
            return ""
        images = convert_from_path(str(pdf_path), first_page=page_idx + 1, last_page=page_idx + 1, dpi=200)
        if not images:
            return ""
        text = pytesseract.image_to_string(images[0], lang="eng+fra")
        return text.strip()
    except Exception:
        return ""


def pdf_merge(input_paths: list[str], output_path: str) -> str:
    """Merge several PDFs into one."""
    try:
        from pypdf import PdfWriter
    except ImportError:
        return _err("pypdf manquant")
    if not input_paths:
        return _err("input_paths vide")
    writer = PdfWriter()
    for ip in input_paths:
        sp = _resolve(ip)
        if not sp.exists():
            return _err(f"introuvable: {sp}")
        writer.append(str(sp))
    out = _resolve(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "wb") as f:
        writer.write(f)
    return _ok(f"PDF merge ({len(input_paths)}) → {out}")


def pdf_split(path: str, output_dir: str, ranges: str = "") -> str:
    """Split a PDF. ranges='1-3,4-6' creates one file per range. Empty = one file per page."""
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        return _err("pypdf manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    out_dir = _resolve(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    reader = PdfReader(str(p))
    total = len(reader.pages)
    if ranges:
        groups = []
        for r in ranges.split(","):
            r = r.strip()
            if "-" in r:
                a, b = r.split("-")
                groups.append(range(int(a) - 1, int(b)))
            else:
                groups.append(range(int(r) - 1, int(r)))
    else:
        groups = [range(i, i + 1) for i in range(total)]
    written = []
    for g in groups:
        w = PdfWriter()
        for i in g:
            if 0 <= i < total:
                w.add_page(reader.pages[i])
        first = list(g)[0] + 1 if list(g) else 0
        last = list(g)[-1] + 1 if list(g) else 0
        name = f"{p.stem}_p{first}-{last}.pdf" if first != last else f"{p.stem}_p{first}.pdf"
        outp = out_dir / name
        with open(outp, "wb") as f:
            w.write(f)
        written.append(str(outp))
    return _ok(f"PDF split → {len(written)} fichiers dans {out_dir}")


def pdf_extract_pages(path: str, pages: str, output_path: str) -> str:
    """Extract specific pages into a new PDF. pages='1,3-5,8'."""
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        return _err("pypdf manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    reader = PdfReader(str(p))
    total = len(reader.pages)
    targets = _parse_pages(pages, total)
    writer = PdfWriter()
    for i in targets:
        if 0 <= i < total:
            writer.add_page(reader.pages[i])
    out = _resolve(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "wb") as f:
        writer.write(f)
    return _ok(f"PDF {len(targets)} pages → {out}")


def pdf_rotate(path: str, output_path: str, angle: int = 90, pages: str = "") -> str:
    """Rotate pages of a PDF (angle: 90, 180, 270)."""
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        return _err("pypdf manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    reader = PdfReader(str(p))
    total = len(reader.pages)
    targets = set(_parse_pages(pages, total)) if pages else set(range(total))
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        if i in targets:
            page.rotate(angle)
        writer.add_page(page)
    out = _resolve(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "wb") as f:
        writer.write(f)
    return _ok(f"PDF rotation {angle}° → {out}")


def pdf_metadata(path: str) -> str:
    """Get PDF metadata (title, author, pages, etc.)."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return _err("pypdf manquant")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    reader = PdfReader(str(p))
    meta = reader.metadata or {}
    info = {
        "pages": len(reader.pages),
        "title": str(meta.get("/Title", "")),
        "author": str(meta.get("/Author", "")),
        "subject": str(meta.get("/Subject", "")),
        "creator": str(meta.get("/Creator", "")),
        "producer": str(meta.get("/Producer", "")),
        "encrypted": reader.is_encrypted,
    }
    return json.dumps(info, ensure_ascii=False, indent=2)


def pdf_add_watermark(path: str, output_path: str, text: str,
                      font_size: int = 50, opacity: float = 0.3) -> str:
    """Add a diagonal text watermark to every page."""
    try:
        from pypdf import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        import io
    except ImportError:
        return _err("pypdf+reportlab manquants")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    # Build watermark PDF in memory
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    c.setFont("Helvetica", font_size)
    c.setFillGray(0.5, opacity)
    c.saveState()
    c.translate(letter[0] / 2, letter[1] / 2)
    c.rotate(45)
    c.drawCentredString(0, 0, text)
    c.restoreState()
    c.save()
    buf.seek(0)
    wm_reader = PdfReader(buf)
    wm_page = wm_reader.pages[0]
    reader = PdfReader(str(p))
    writer = PdfWriter()
    for page in reader.pages:
        page.merge_page(wm_page)
        writer.add_page(page)
    out = _resolve(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "wb") as f:
        writer.write(f)
    return _ok(f"PDF watermark → {out}")


def pdf_encrypt(path: str, output_path: str, password: str,
                owner_password: str = "") -> str:
    """Password-protect a PDF."""
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        return _err("pypdf manquant")
    if not password:
        return _err("password requis")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    reader = PdfReader(str(p))
    writer = PdfWriter(clone_from=reader)
    writer.encrypt(user_password=password, owner_password=owner_password or password)
    out = _resolve(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "wb") as f:
        writer.write(f)
    return _ok(f"PDF chiffré → {out}")


def _parse_pages(spec: str, total: int) -> list[int]:
    """Parse '1,3-5,8' → [0, 2, 3, 4, 7] (0-indexed)."""
    out = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-")
            out.extend(range(int(a) - 1, int(b)))
        else:
            out.append(int(part) - 1)
    return out


# ─── Email (.eml) ────────────────────────────────────────────────────────────

def eml_create(path: str, to: str, subject: str, body: str,
               from_addr: str = "", cc: str = "", attachments: list[str] | None = None) -> str:
    """Create a standard .eml file (importable in any mail client)."""
    try:
        from email.message import EmailMessage
        import mimetypes
    except ImportError:
        return _err("stdlib email indisponible")
    msg = EmailMessage()
    msg["To"] = to
    msg["Subject"] = subject
    if from_addr:
        msg["From"] = from_addr
    if cc:
        msg["Cc"] = cc
    msg.set_content(body)
    for att in attachments or []:
        ap = _resolve(att)
        if not ap.exists():
            continue
        ctype, _ = mimetypes.guess_type(str(ap))
        maintype, _, subtype = (ctype or "application/octet-stream").partition("/")
        msg.add_attachment(ap.read_bytes(), maintype=maintype, subtype=subtype, filename=ap.name)
    p = _resolve(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_bytes(bytes(msg))
    return _ok(f"eml → {p}")


def eml_read(path: str) -> str:
    """Parse an .eml file, return JSON {from, to, subject, body, attachments}."""
    try:
        from email import policy
        from email.parser import BytesParser
    except ImportError:
        return _err("stdlib email indisponible")
    p = _resolve(path)
    if not p.exists():
        return _err(f"introuvable: {p}")
    with open(p, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)
    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body = part.get_content()
                break
    else:
        body = msg.get_content()
    attachments = []
    for part in msg.iter_attachments():
        attachments.append({"filename": part.get_filename(), "type": part.get_content_type()})
    out = {
        "from": str(msg.get("From", "")),
        "to": str(msg.get("To", "")),
        "cc": str(msg.get("Cc", "")),
        "subject": str(msg.get("Subject", "")),
        "date": str(msg.get("Date", "")),
        "body": body[:8000],
        "attachments": attachments,
    }
    return json.dumps(out, ensure_ascii=False, indent=2)
